#!/usr/bin/env -S uv run --quiet --script
# /// script
# requires-python = ">=3.10"
# dependencies = [
#   "pyyaml>=6.0",
#   "python-pptx>=0.6.21",
# ]
# ///
"""
export_pptx.py — 将 slides.yml + images/ 导出为 .pptx 文件

每张幻灯片图片作为全屏背景填满 16:9 页面，幻灯片标题写入备注栏（方便后续编辑）。

Usage:
  uv run scripts/export_pptx.py --slides-yml output/ppt-gen/slides.yml
  uv run scripts/export_pptx.py --slides-yml output/ppt-gen/slides.yml --output output/ppt-gen/presentation.pptx
  uv run scripts/export_pptx.py --slides-yml output/ppt-gen/slides.yml --output deck.pptx --title "我的演示文稿"
"""

import argparse
import sys
from pathlib import Path

import yaml
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


# 16:9 标准尺寸（单位：EMU）
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


def load_slides_yml(yml_path: str) -> dict:
    path = Path(yml_path)
    if not path.exists():
        print(f"Error: slides.yml not found at {yml_path}")
        sys.exit(1)
    with open(path, "r", encoding="utf-8") as f:
        data = yaml.safe_load(f)
    presentation = data.get("presentation", data)
    return presentation


def build_pptx(slides_yml_path: str, output_path: str, override_title: str = None) -> str:
    """
    Build a .pptx from slides.yml + images/.
    Returns the output path.
    """
    base_dir = Path(slides_yml_path).parent
    data = load_slides_yml(slides_yml_path)

    title = override_title or data.get("title", "Presentation")
    slides = data.get("slides", [])
    slides = sorted(slides, key=lambda s: s.get("index", 9999))

    if not slides:
        print("Warning: No slides found in slides.yml")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    blank_layout = prs.slide_layouts[6]  # completely blank layout

    success = 0
    skipped = 0

    for slide_data in slides:
        idx = slide_data.get("index", success + skipped + 1)
        nn = str(idx).zfill(2)
        slide_title = slide_data.get("title", f"Slide {idx}")
        subtitle = slide_data.get("subtitle", "")
        key_message = slide_data.get("key_message", "")

        # Resolve image path
        image_output = slide_data.get("image", {}).get("output", f"images/slide-{nn}.jpg")
        # 如果 image_output 已是绝对路径，直接用；否则相对 slides.yml 所在目录拼接
        image_path = Path(image_output) if Path(image_output).is_absolute() else base_dir / image_output

        # Fallback: try common extensions
        if not image_path.exists():
            for ext in [".jpg", ".jpeg", ".png", ".webp"]:
                candidate = image_path.with_suffix(ext)
                if candidate.exists():
                    image_path = candidate
                    break

        slide = prs.slides.add_slide(blank_layout)

        if image_path.exists():
            # Add image as full-bleed background
            pic = slide.shapes.add_picture(
                str(image_path),
                left=Emu(0),
                top=Emu(0),
                width=SLIDE_WIDTH,
                height=SLIDE_HEIGHT,
            )
            # Push image to back
            slide.shapes._spTree.remove(pic._element)
            slide.shapes._spTree.insert(2, pic._element)
            success += 1
            print(f"  ✅ [{idx:02d}] {slide_title} — {image_path.name}")
        else:
            # No image found: add a dark placeholder with title text
            from pptx.util import Pt
            txBox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_title
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(36)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # Dark background rectangle
            from pptx.util import Emu as _Emu
            bg = slide.shapes.add_shape(
                1,  # MSO_SHAPE_TYPE.RECTANGLE
                _Emu(0), _Emu(0), SLIDE_WIDTH, SLIDE_HEIGHT
            )
            bg.fill.solid()
            bg.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
            bg.line.fill.background()
            slide.shapes._spTree.remove(bg._element)
            slide.shapes._spTree.insert(2, bg._element)

            skipped += 1
            print(f"  ⚠️  [{idx:02d}] {slide_title} — image not found ({image_output}), using text placeholder")

        # Write title + key_message into notes (useful for later editing)
        notes_slide = slide.notes_slide
        notes_tf = notes_slide.notes_text_frame
        notes_lines = [slide_title]
        if subtitle:
            notes_lines.append(subtitle)
        if key_message:
            notes_lines.append(key_message)
        notes_tf.text = "\n".join(notes_lines)

    # Save
    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))

    print(f"\n✅ PPTX saved: {out}")
    print(f"   Title  : {title}")
    print(f"   Slides : {success + skipped} total ({success} with images, {skipped} placeholders)")
    return str(out)


def main() -> int:
    parser = argparse.ArgumentParser(
        description="将 slides.yml + images/ 导出为 .pptx",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  python scripts/export_pptx.py --slides-yml output/ppt-gen/slides.yml
  python scripts/export_pptx.py --slides-yml output/ppt-gen/slides.yml --output deck.pptx
  python scripts/export_pptx.py --slides-yml output/ppt-gen/slides.yml --output deck.pptx --title "Q2 Brand Launch"

Notes:
  - Each slide image is placed as a full-bleed 16:9 background.
  - Slide title and key_message are written to the notes pane for reference.
  - If an image file is missing, a dark text-placeholder slide is inserted instead.
  - Requires: pip install python-pptx pyyaml
""",
    )
    parser.add_argument("--slides-yml", required=True, help="Path to slides.yml")
    parser.add_argument(
        "--output", "-o", default=None,
        help="Output .pptx path (default: same directory as slides.yml, named after presentation title)"
    )
    parser.add_argument("--title", default=None, help="Override presentation title")

    args = parser.parse_args()

    # Default output path: same dir as slides.yml
    if args.output is None:
        yml_path = Path(args.slides_yml)
        data = load_slides_yml(args.slides_yml)
        pptx_title = args.title or data.get("title", "presentation")
        safe_title = "".join(c if c.isalnum() or c in " -_" else "_" for c in pptx_title).strip()
        args.output = str(yml_path.parent / f"{safe_title}.pptx")

    build_pptx(args.slides_yml, args.output, args.title)
    return 0


if __name__ == "__main__":
    sys.exit(main())
